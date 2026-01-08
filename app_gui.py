import os
import shutil
import json
import re
import threading
from pathlib import Path
import ollama
import customtkinter as ctk
from tkinter import filedialog
import PyPDF2
from pptx import Presentation
from docx import Document
from datetime import datetime

# GUI
ctk.set_appearance_mode("Dark")
ctk.set_default_color_theme("dark-blue")

# --- CONFIG ---
MODEL_NAME = "llama3"

# CATEGORY MAP
STRUCTURE_MAP = """
- Career/CV
- Career/Contracts
- Career/Documents
- Programming/Python
- Programming/WebDev
- Programming/C_CPP
- Programming/DataScience
- University/Homework
- University/Projects
- University/Materials
- Media/Images
- Media/Videos
- Others
"""

# --- LOGIC ---

def organize_media_by_date(file_path):
    """
    BRAKUJĄCA FUNKCJA: Sortuje media po dacie (Rok/Miesiąc).
    """
    ext = file_path.suffix.lower()
    
    # CATEGORY 
    if ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.heic', '.webp']:
        base_cat = "Media/Images"
    elif ext in ['.mp4', '.mov', '.avi', '.mkv', '.webm']:
        base_cat = "Media/Videos"
    else:
        return None, None

    try:
        # looking for file infos
        timestamp = os.path.getmtime(file_path)
        date = datetime.fromtimestamp(timestamp)
        
        year = date.strftime("%Y")
        month = date.strftime("%m_%B") # example: "05_May"
        
        # result: "Media/Images/2023/05_May"
        category_path = f"{base_cat}/{year}/{month}"
        
        return category_path, file_path.name
    except Exception:
        return None, None

def analyze_file(filename, snippet):
    """Ask Ollama for Category and filename"""
    prompt = f"""
    You are an intelligent file management agent.
    Your goal is to organize files based on their content and filename.
    
    CATEGORY MAP:
    {STRUCTURE_MAP}
    
    FILE TO ANALYZE:
    - Original Name: "{filename}"
    - Content Snippet: "{snippet}"
    
    INSTRUCTIONS:
    1. Analyze the content and name to pick the BEST category from the Map.
    2. Generate a new, clean filename in English (use snake_case, e.g., 'invoice_march_2024.pdf').
    3. Keep the original file extension.
    4. If the content is unclear, use 'Others'.
    
    OUTPUT FORMAT:
    Return ONLY a raw JSON object. Do not write markdown blocks like ```json.
    Example:
    {{
        "category": "Programming/Python",
        "new_name": "sorting_algorithm.py"
    }}
    """
    
    try:
        response = ollama.chat(
            model=MODEL_NAME, 
            messages=[{'role': 'user', 'content': prompt}],
            format='json',
            options={'temperature': 0.1}
        )
        response_content = response['message']['content']
        return json.loads(response_content)

    except Exception as e:
        print(f"Błąd AI dla {filename}: {e}")
        return {"category": "Others", "new_name": filename}

def extract_text_snippet(file_path):
    """read from file"""
    try:
        ext = file_path.suffix.lower()
        if ext in ['.txt', '.md', '.py', '.js', '.html', '.css', '.cpp', '.java']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read(800)
        elif ext == '.pdf':
            reader = PyPDF2.PdfReader(file_path)
            if len(reader.pages) > 0: return reader.pages[0].extract_text()[:800]
        elif ext == '.docx':
            doc = Document(file_path)
            full_text = []
            for para in doc.paragraphs:
                if para.text.strip(): full_text.append(para.text)
                if len(full_text) > 20: break
            return "\n".join(full_text)[:800]
        elif ext == '.pptx':
            prs = Presentation(file_path)
            text_runs = []
            for slide in prs.slides[:2]:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): text_runs.append(shape.text)
            return " ".join(text_runs)[:800]
    except: return ""
    return ""

def get_unique_path(path):
    """avoid file overwrite"""
    if not path.exists(): return path
    counter = 1
    while True:
        new_path = path.parent / f"{path.stem}_{counter}{path.suffix}"
        if not new_path.exists(): return new_path
        counter += 1

def sanitize_filename(name):
    return re.sub(r'[^\w\-_.]', '_', name)

# --- GUI ---

class SmartSortApp(ctk.CTk):
    def __init__(self):
        super().__init__()
        self.title("DocFlow")
        self.geometry("800x600")
        
        # Nagłówek
        self.header = ctk.CTkLabel(self, text="DocFlow", font=("Roboto", 24, "bold"))
        self.header.pack(pady=20)

        # Sekcja folderu
        self.frame_input = ctk.CTkFrame(self)
        self.frame_input.pack(pady=10, padx=20, fill="x")

        self.path_entry = ctk.CTkEntry(self.frame_input, placeholder_text="Select folder...", height=35)
        self.path_entry.pack(side="left", fill="x", expand=True, padx=10, pady=10)
        
        self.browse_btn = ctk.CTkButton(self.frame_input, text="Browse", command=self.browse_folder)
        self.browse_btn.pack(side="right", padx=10)

        # Przycisk Start
        self.start_btn = ctk.CTkButton(self, text="START SORTING", command=self.start_process, 
                                       fg_color="#2CC985", hover_color="#229966", height=40)
        self.start_btn.pack(pady=20)

        # Logi
        self.log_box = ctk.CTkTextbox(self, width=700, height=300, font=("Consolas", 12))
        self.log_box.pack(pady=10)
        self.log_box.insert("0.0", ">> Ready. Ensure Ollama is running in the background.\n")

    def browse_folder(self):
        folder = filedialog.askdirectory()
        if folder:
            self.path_entry.delete(0, "end")
            self.path_entry.insert(0, folder)

    def log(self, message):
        self.log_box.insert("end", message + "\n")
        self.log_box.see("end")

    def start_process(self):
        source = self.path_entry.get()
        if not source:
            self.log("!! ERROR: Please select a folder first!   ")
            return
        
        self.start_btn.configure(state="disabled", text="Sorting...", fg_color="gray")
        threading.Thread(target=self.run_sorting, args=(source,)).start()
    
    def run_sorting(self, source_folder):
        source_path = Path(source_folder)
        target_path = source_path / "Sorted_AI"
        
        files = [f for f in source_path.iterdir() if f.is_file()]
        self.log(f">> Found {len(files)} files. Starting...")

        for file in files:
            try:
                ext = file.suffix.lower()
                self.log(f"Processing: {file.name}...")

                # 1. ARCHIVES
                if ext in ['.zip', '.rar', '.7z', '.tar', '.gz', '.pkg']:
                    final_folder = target_path / "Archives"
                    final_folder.mkdir(parents=True, exist_ok=True)
                    dest_path = get_unique_path(final_folder / file.name)
                    shutil.move(str(file), str(dest_path))
                    self.log(f"📚 [Archives] -> {dest_path.name}")
                    continue 

                # 2. APPS
                if ext in ['.exe', '.msi', '.bat', '.sh', '.iso', '.dmg']:
                    final_folder = target_path / "Apps"
                    final_folder.mkdir(parents=True, exist_ok=True)
                    dest_path = get_unique_path(final_folder / file.name)
                    shutil.move(str(file), str(dest_path))
                    self.log(f"📦 [Apps] -> {dest_path.name}")
                    continue 

                # 3. MEDIA
                media_cat, media_name = organize_media_by_date(file)
                if media_cat:
                    final_folder = target_path / media_cat
                    final_folder.mkdir(parents=True, exist_ok=True)
                    dest_path = get_unique_path(final_folder / media_name)
                    shutil.move(str(file), str(dest_path))
                    self.log(f"📸 [{media_cat}] -> {dest_path.name}")
                    continue

                # 4. TECH FILES
                misc_extensions = ['.yml', '.yaml', '.ini', '.cfg', '.log', '.tmp', '.bak', '.config', '.xml']
                if ext in misc_extensions:
                    final_folder = target_path / "Others"
                    final_folder.mkdir(parents=True, exist_ok=True)
                    dest_path = get_unique_path(final_folder / file.name)
                    shutil.move(str(file), str(dest_path))
                    self.log(f"⚙️ [Others/Config] -> {dest_path.name}")
                    continue

                # 5. AI ANALYSER
                snippet = extract_text_snippet(file)
                result = analyze_file(file.name, snippet)
                
                cat = result.get("category", "Others")
                new_name = result.get("new_name", file.name)

                if "Media" in cat: cat = "Others"
                
                final_folder = target_path / cat
                final_folder.mkdir(parents=True, exist_ok=True)
                
                safe_name = sanitize_filename(Path(new_name).stem) + file.suffix
                dest_path = get_unique_path(final_folder / safe_name)
                
                shutil.move(str(file), str(dest_path))
                self.log(f"🧠 [{cat}] -> {dest_path.name}")
                
            except Exception as e:
                self.log(f"❌ Error with {file.name}: {e}")

        self.log(">> DONE!")
        self.start_btn.configure(state="normal", text="RUN SORTING", fg_color="green")

if __name__ == "__main__":
    app = SmartSortApp()
    app.mainloop()
